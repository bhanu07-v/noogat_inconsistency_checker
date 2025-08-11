import os
import re
import json
import textwrap
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import pytesseract
from dotenv import load_dotenv
from rapidfuzz import fuzz
import google.generativeai as genai
from dateutil import parser as dateparser

# Load API key from .env
load_dotenv()
API_KEY = os.getenv("GEMINI_API_KEY")
if not API_KEY:
    raise ValueError("Gemini API key not found. Please set GEMINI_API_KEY in .env")

# Configure Gemini API
genai.configure(api_key=API_KEY)

def extract_text_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    slides = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text and shape.text.strip():
                texts.append(shape.text.strip())
        slides.append({"slide": i, "text": "\n".join(texts)})
    return slides, prs

def extract_images_from_pptx(prs, out_dir):
    out_dir.mkdir(parents=True, exist_ok=True)
    saved = []
    for i, slide in enumerate(prs.slides, start=1):
        for j, shape in enumerate(slide.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img = shape.image
                ext = img.ext
                filename = f"slide_{i}_img_{j}.{ext}"
                path = out_dir / filename
                with open(path, "wb") as f:
                    f.write(img.blob)
                saved.append((i, path))
    return saved

def ocr_image(path):
    return pytesseract.image_to_string(Image.open(path))

def normalize_number(raw):
    s = raw.strip().lower().replace(",", "")
    multiplier = 1.0
    if "thousand" in s or s.endswith("k"):
        multiplier = 1e3
    if "million" in s or s.endswith("m"):
        multiplier = 1e6
    if "billion" in s or s.endswith("b"):
        multiplier = 1e9
    s2 = re.sub(r'[^\d\.]', '', s)
    try:
        return float(s2) * multiplier
    except:
        return None

def similar(a, b):
    return fuzz.token_sort_ratio(a, b) / 100.0

def local_numeric_conflicts(mentions):
    issues = []
    nums = [m for m in mentions if m["type"] == "number" and m["value"] is not None]
    for i in range(len(nums)):
        for j in range(i+1, len(nums)):
            a, b = nums[i], nums[j]
            if a["slide"] == b["slide"]:
                continue
            if similar(a["context"], b["context"]) > 0.45:
                rel = abs(a["value"] - b["value"]) / max(abs(a["value"]), abs(b["value"]))
                if rel > 0.10:
                    issues.append({
                        "type": "number_conflict",
                        "slides": [a["slide"], b["slide"]],
                        "a_raw": a["raw"], "b_raw": b["raw"],
                        "context_a": a["context"], "context_b": b["context"]
                    })
    return issues

def extract_mentions(slides):
    NUM_RE = re.compile(r'(\$|₹)?\s*([0-9][0-9,\.]*\s*(?:K|M|B|million|billion|thousand)?)', re.I)
    PCT_RE = re.compile(r'(\d+(?:\.\d+)?)\s*%')
    DATE_RE = re.compile(r'\b(?:Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec|\d{1,2})(?:[^\n,]{0,20})\d{2,4}\b', re.I)
    mentions = []
    for s in slides:
        text = s["text"]
        for m in NUM_RE.finditer(text):
            raw = m.group(0).strip()
            val = normalize_number(raw)
            ctx = text[max(0, m.start()-50):m.end()+50].replace("\n"," ")
            mentions.append({"slide": s["slide"], "raw": raw, "value": val, "type": "number", "context": ctx})
        for m in PCT_RE.finditer(text):
            raw = m.group(0)
            try: val = float(m.group(1))/100.0
            except: val = None
            ctx = text[max(0, m.start()-50):m.end()+50].replace("\n"," ")
            mentions.append({"slide": s["slide"], "raw": raw, "value": val, "type": "percent", "context": ctx})
        for m in DATE_RE.finditer(text):
            raw = m.group(0)
            try: parsed = dateparser.parse(raw).isoformat()
            except: parsed = None
            ctx = text[max(0, m.start()-50):m.end()+50].replace("\n"," ")
            mentions.append({"slide": s["slide"], "raw": raw, "value": parsed, "type": "date", "context": ctx})
    return mentions

def call_gemini_for_inconsistencies(slides):
    prompt = (
        "You are an AI that finds factual or logical inconsistencies in presentation slides.\n"
        "Compare the slides and find:\n"
        "- Conflicting numbers\n"
        "- Contradictory statements\n"
        "- Timeline mismatches\n"
        "Return ONLY JSON in format:\n"
        '[{"type": "...", "slides": [..], "summary": "...", "evidence": ["...", "..."]}]'
    )
    resp = genai.GenerativeModel("gemini-2.0-flash").generate_content(
        prompt + "\nSlides:\n" + json.dumps(slides, indent=2)
    )
    try:
        return json.loads(resp.text)
    except:
        return []

def main():
    pptx_path = input("Enter path to PPTX file: ").strip('"')
    out_dir = Path("output")
    out_dir.mkdir(exist_ok=True)

    slides, prs = extract_text_from_pptx(pptx_path)
    images = extract_images_from_pptx(prs, out_dir / "images")

    # OCR for images
    for slide_num, img_path in images:
        ocr_text = ocr_image(img_path)
        if ocr_text.strip():
            slides[slide_num-1]["text"] += "\n[OCR]\n" + ocr_text.strip()

    mentions = extract_mentions(slides)
    local_issues = local_numeric_conflicts(mentions)

    gemini_issues = call_gemini_for_inconsistencies(slides)

    final_issues = {"local": local_issues, "gemini": gemini_issues}

    # Save JSON
    with open(out_dir / "inconsistencies.json", "w", encoding="utf-8") as f:
        json.dump(final_issues, f, indent=2, ensure_ascii=False)

    # Save human-readable report
    lines = ["Inconsistency Report\n"]
    for issue in local_issues:
        lines.append(f"[LOCAL {issue['type']}] Slides {issue['slides']} — {issue['a_raw']} vs {issue['b_raw']}")
    for issue in gemini_issues:
        lines.append(f"[GEMINI {issue['type']}] Slides {issue['slides']} — {issue['summary']}")
    with open(out_dir / "report.txt", "w", encoding="utf-8") as f:
        f.write("\n".join(lines))

    print("Analysis complete. See 'output/inconsistencies.json' and 'output/report.txt'.")

if __name__ == "__main__":
    main()
